#!/usr/bin/env python3
"""Generate Professional BLIP-2 Reproduction Presentation"""

import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    print("✓ python-pptx imported successfully")
except ImportError as e:
    print(f"Error: {e}")
    print("Installing python-pptx...")
    import subprocess
    subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

def create_professional_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    DARK_BG = RGBColor(15, 23, 42)  
    TITLE_COLOR = RGBColor(59, 130, 246)  
    TEXT_COLOR = RGBColor(255, 255, 255)  
    DARK_TEXT = RGBColor(20, 20, 30)  
    ACCENT_COLOR = RGBColor(34, 197, 94)  
    CARD_BG = RGBColor(30, 41, 59)  
    
    def add_slide_with_bg():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        return slide
    
    def add_title_bar(slide, title):
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = TITLE_COLOR
        title_shape.line.fill.background()
        
        text_frame = title_shape.text_frame
        text_frame.text = title
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.paragraph_format.left_indent = Pt(20)
    
    def add_two_column(slide, left_title, left_items, right_title, right_items):
        left_header = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(4.5), Inches(0.35))
        tf = left_header.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        left_content = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(4.5), Inches(5.7))
        tf = left_content.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            p.line_spacing = 1.15
            if item.startswith('•'):
                p.paragraph_format.left_indent = Pt(15)
        
        right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.35))
        tf = right_header.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        right_content = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.5), Inches(5.7))
        tf = right_content.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_items):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            p.line_spacing = 1.15
            if item.startswith('•'):
                p.paragraph_format.left_indent = Pt(15)
    
    # SLIDE 1: TITLE
    slide = add_slide_with_bg()
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Multimodal LLMs\nFrom Theory to Practice"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(1.5))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Comparative Critique + BLIP-2 Local Reproduction"
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # SLIDE 2-15: Content slides (simplified for space)
    slide_configs = [
        ("Three Papers: Quick Comparison", [
            "BLIP-2", ["Q-Former bridge", "129M pairs", "43.7 BLEU-4", "✓ Best target"],
            "LLaVA & VisionLLM2", ["Instruction-tuned", "158K pairs", "92.53% ScienceQA", "⚠ Constraints"]
        ]),
        ("Setup: Local Downscaling", [
            "Hardware & Models", ["RTX 3070 8GB", "CLIP ViT-L", "OPT-350M", "224px res"],
            "Paper Configuration", ["16×A100 40GB", "ViT-g", "OPT-2.7B", "364px res"]
        ]),
        ("Local Results Scaling", [
            "10k Pilot", ["BLEU-4: 1.45", "CIDEr: 3.03", "Baseline"],
            "50k Main", ["BLEU-4: 11.13", "CIDEr: 37.57", "✓ PEAK"]
        ]),
        ("Performance Gap Analysis", [
            "Sources of Gap", ["Backbone: 3-4×", "LLM: 8×", "Data: 1290×", "Resolution/budget"],
            "Gap Interpretation", ["~25-30× total", "Multiplicative", "NOT pipeline failure", "Scale-driven"]
        ]),
        ("100k Run: Peak-Early Phenomenon", [
            "Training Stable", ["Losses smooth", "No divergence", "Stage 1: 0.33", "Stage 2: 0.36"],
            "Metrics Degraded", ["Best: epoch 2", "Final: epoch 14", "-23% BLEU-4", "Overfitting"]
        ]),
        ("Main Conclusions", [
            "Pipeline Status", ["✓ Reproducible", "All 3 stages work", "On RTX 3070", "End-to-end OK"],
            "Performance Status", ["✗ Not reproduced", "25-30× gap", "Scale-driven", "Architecture sound"]
        ]),
        ("Key Takeaways", [
            "Technical", ["Backbone critical", "Selection > schedules", "Frozen design works", "Modular+scalable"],
            "Reproducibility", ["Pipeline ≠ Performance", "Local tuning needed", "Artifact trail matters", "Scale compounds"]
        ]),
        ("Reproducibility Trail", [
            "Artifacts Available", ["Checkpoint registry", "Per-epoch evals", "Config files", "Loss logs"],
            "Enables Verification", ["Full trajectory", "Model selection audit", "Hyperparams exact", "Training stable"]
        ]),
    ]
    
    for title, config_data in slide_configs:
        if len(config_data) == 4:
            slide = add_slide_with_bg()
            add_title_bar(slide, title)
            add_two_column(slide, config_data[0], config_data[1], config_data[2], config_data[3])
    
    # Additional summary slide
    slide = add_slide_with_bg()
    add_title_bar(slide, "Research Summary")
    
    summary_text = slide.shapes.add_textbox(Inches(0.4), Inches(0.85), Inches(9.2), Inches(6.4))
    tf = summary_text.text_frame
    tf.word_wrap = True
    
    summary = [
        "Study evaluated three landmark multimodal LLM papers (BLIP-2, LLaVA, VisionLLM v2)",
        "through architectural critique and student-scale reproducibility lens.",
        "",
        "BLIP-2 is the most defensible local target. Its frozen-backbone modularity provides",
        "coherent downscaling path. Reproduced all 3 stages on RTX 3070 8GB successfully.",
        "",
        "Best result: BLEU-4 11.13, CIDEr 37.57 (50k subset) vs. 43.7/145.8 (paper).",
        "25-30× gap is systematic (backbones+data+compute), not pathological.",
        ""," For practitioners: Invest in backbone understanding, use metric-driven selection.",
        "For researchers: Pipeline ≠ performance reproducibility (crucial distinction)."
    ]
    
    for i, line in enumerate(summary):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = ACCENT_COLOR if any(s in line for s in ["BLIP-2", "●"]) else TEXT_COLOR
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        p.line_spacing = 1.2
    
    # Save
    prs.save("BLIP2_Reproduction_Presentation_Professional.pptx")
    print("✓ Professional presentation created!")
    print("✓ File: BLIP2_Reproduction_Presentation_Professional.pptx")

if __name__ == "__main__":
    create_professional_presentation()
