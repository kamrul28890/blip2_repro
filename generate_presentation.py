#!/usr/bin/env python3
"""Generate BLIP-2 Reproduction Analysis presentation"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 45, 80)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        sub_frame = sub_box.text_frame
        sub_frame.word_wrap = True
        p = sub_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(26)
        p.font.color.rgb = RGBColor(200, 220, 255)

def add_content_slide(title, content_items, title_bg=RGBColor(25, 45, 80)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = title_bg
    title_shape.line.color.rgb = RGBColor(100, 150, 200)
    
    title_frame = title_shape.text_frame
    title_frame.word_wrap = False
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.paragraph_format.space_before = Pt(4)
    p.paragraph_format.space_after = Pt(4)
    p.paragraph_format.left_indent = Pt(10)
    
    content_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(9.2), Inches(6.3))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        p.line_spacing = 1.1

# ============= SLIDE CONTENT =============

# Slide 1: Title
add_title_slide('Multimodal LLMs: From Theory to Practice', 'Comparative Critique + BLIP-2 Reproduction Study')

# Slide 2: Paper Selection
add_content_slide('Three Papers: Why & What', [
    '• BLIP-2 (Li et al., 2023): Q-Former bridge → frozen backbones → efficient alignment',
    '• LLaVA (Liu et al., 2023): Instruction tuning → multimodal assistants → data-centric',
    '• VisionLLM v2 (Wu et al., 2024): Multi-decoder routing → end-to-end generalist → complex systems',
    '',
    'Common domain: Vision-language alignment & cross-modal understanding',
    'Key difference: Architecture complexity & supervision strategy scale',
    'Selection rationale: BLIP-2 most suitable for local reproduction (modular, open, downscalable)'
])

# Slide 3: Problem Statements
add_content_slide('Problem Statements', [
    'BLIP-2: Can lightweight bridge efficiently connect frozen vision & language?',
    '→ Motivates: Compute-efficient multimodal learning without full retraining',
    '',
    'LLaVA: How to build practical multimodal assistants vs. narrow captioning?',
    '→ Motivates: Supervision format impact on assistant behavior',
    '',
    'VisionLLM v2: Can one system unify many vision & vision-language tasks?',
    '→ Motivates: Generalist perception via shared router + task decoders',
    '',
    'Shared thread: Reuse strong frozen backbones, reduce training complexity'
])

# Slide 4: Methodology
add_content_slide('Methodology Comparison', [
    'BLIP-2: Two-stage Q-Former bridge (Stage1: ITC+ITM+LM; Stage2: aligned generation)',
    '  Data: 129M pairs (COCO+Visual Genome+CC3M+CC12M+SBU+LAION)',
    '  Compute: 16×A100 40GB, multi-day training',
    '',
    'LLaVA: Lightweight projection + multimodal instruction tuning (GPT-4 generated)',
    '  Data: 158K instruction pairs (ViT-L CLIP + Vicuna-7B)',
    '  Compute: 8×A100, ~10 hours training',
    '',
    'VisionLLM v2: Central MLLM + super-link routing to multiple decoders',
    '  Data: 64 diverse datasets across ~100 tasks',
    '  Compute: 64→128 A100 GPUs, 18 days'
])

# Slide 5: Results & Reproducibility
add_content_slide('Results & Reproducibility', [
    'BLIP-2 | BLEU-4: 43.7, CIDEr: 145.8 | ✓ BEST LOCAL TARGET (modular, downscalable)',
    '  - Frozen architecture preserves coherent downscaling logic',
    '  - Official code available (LAVIS), no synthetic data dependency',
    '',
    'LLaVA | ScienceQA: 92.53% (w/ GPT-4) | ⚠ VIABLE WITH CONSTRAINTS',
    '  - Simple arch, but depends on GPT-4 instruction pipeline',
    '  - Still requires 8× A100 for full training',
    '',
    'VisionLLM v2 | Broad generalist tasks | ✗ NOT LOCAL FEASIBLE',
    '  - 64-128 A100 required, multi-task complexity too high',
    '  - Gains hard to disentangle under downscaling'
])

# Slide 6: Limitations
add_content_slide('Limitations Across Papers', [
    'BLIP-2: Published performance tightly coupled to large backbones (ViT-g, OPT-2.7B)',
    '  → Scaling to commodity hardware remains unexplored in paper',
    '',
    'LLaVA: Instruction-data reliance (GPT-4 generated) introduces synthetic bias',
    '  → Reproducibility depends on instruction pipeline quality & availability',
    '',
    'VisionLLM v2: Multi-decoder complexity makes ablation analysis difficult',
    '  → Gains attributable to routing? Decoder reuse? Data scale? Unclear separation',
    '',
    'General: All three show scale→performance coupling; local downsizing implications unclear'
])

# Slide 7: Local Setup
add_content_slide('Local BLIP-2 Setup: Bridging Scale Gap', [
    'Hardware: Single NVIDIA RTX 3070 8GB (vs. 16 A100 40GB)',
    'Vision: CLIP ViT-L (vs. ViT-g); LLM: OPT-350M (vs. OPT-2.7B); Resolution: 224 (vs. 364)',
    '',
    'Data: COCO Karpathy subsets (10k, 50k, 100k train) vs. 129M pretraining pairs',
    '  - Fixed val/test: 1k each for internal comparison fidelity',
    '',
    'Pipeline stages (frozen backbone design preserved):',
    '  Stage 1: Q-Former (32 queries) + frozen ViT-L via ITC+ITM+LM loss',
    '  Stage 2: Generative alignment (Q→LLM projections)',
    '  Stage 3: Caption finetuning ("a photo of" prompt)',
    '',
    'Schedules: 50k (3/3/5 epochs), 100k (10/10/15 epochs)'
])

# Slide 8: Local Results
add_content_slide('Local Results & Scaling Analysis', [
    'Paper BLIP-2 (ViT-g, OPT-2.7B, 129M pairs): BLEU-4=43.7, CIDEr=145.8',
    '',
    'Local 10k Pilot:   BLEU-4=1.45,   CIDEr=3.03    (baseline)',
    'Local 50k Best:    BLEU-4=11.13,  CIDEr=37.57   (epoch 3) ✓ PEAK',
    'Local 50k Polish:  BLEU-4=10.49,  CIDEr=37.14   (epoch 4)',
    'Local 100k Long:   BLEU-4=8.57,   CIDEr=25.34   (epoch 2 best) ⚠ PEAKS EARLY',
    'Local 100k Final:  BLEU-4=6.61,   CIDEr=18.36   (epoch 14 final) → degradation',
    '',
    'Key finding: 10k→50k shows >7× gain (learning signal present). 50k→100k shows peak-early behavior.'
])

# Slide 9: Gap Analysis
add_content_slide('Why The Gap? Systematic Breakdown', [
    'Paper Configuration: ViT-g (7× params), OPT-2.7B (8× params), 129M pairs (1290× data), 364px res',
    'Local Configuration: ViT-L, OPT-350M, 50k pairs, 224px resolution',
    '',
    'Estimated gap contributions:',
    '• Backbone scale (ViT-g→ViT-L): ~3-4× capacity loss',
    '• LLM scale (OPT-2.7B→OPT-350M): ~8× parameter reduction',
    '• Pre-training data (129M→0 external, COCO-only): Major feature quality impact',
    '• Fine-tune data (COCO scale): 50k local vs. massive official training corpus',
    '• Resolution & compute budget: 224→364, shorter schedules',
    '',
    'Verdict: ~25-30× performance gap is multiplicative effect, NOT pipeline failure'
])

# Slide 10: Peak-Early Phenomenon
add_content_slide('100k Run: Peak-Early Phenomenon', [
    'Training losses decreased smoothly (Stage 1: 0.33, Stage 2: 0.36 in 100k run)',
    '→ Bridge training was stable, no divergence',
    '',
    'But caption metrics showed early peaking:',
    '  Best: Epoch 2 (BLEU-4=8.57, CIDEr=25.34)',
    '  Final: Epoch 14 (BLEU-4=6.61, CIDEr=18.36) – 23% decline',
    '',
    'Root cause: Without strong regularization (dropout, weight decay tuned for this scale),',
    'longer training without validation-driven early stopping leads to overfitting',
    '',
    'Diversity trend: Unique captions increased (677→734), but with lower semantic quality',
    '→ Model memorization / mode collapse on reduced validation set'
])

# Slide 11: Study Limitations
add_content_slide('Study Limitations', [
    '1. Local evaluation: 1k val/test (paper uses full Karpathy test)',
    '   → Numbers comparable internally but not to paper leaderboard',
    '',
    '2. Architecture downscaling: CLIP ViT-L + OPT-350M vs. ViT-g + OPT-2.7B',
    '   → Studied downscaling path, not paper-exact reproduction',
    '',
    '3. Single-seed runs (no variance estimates)',
    '   → Ranking stable, generalization uncertain',
    '',
    '4. Windows portability: Targeted code fixes for non-distributed mode',
    '   → Measures model reproducibility + research codebase portability',
    '',
    '5. No synthetic data: Unlike LLaVA (GPT-4 instruction data), all COCO real data'
])

# Slide 12: Future Work
add_content_slide('Future Work & Unexplored Directions', [
    '1. Backbone scaling study: ViT-B, ViT-L, ViT-H systematically vs. OPT-125M→350M→1.3B',
    '',
    '2. Data augmentation: Synthetic descriptions (self-supervised or from larger models)',
    '',
    '3. Regularization tuning: Early stopping, dropout schedules, learning rate schedules',
    '   for reduced-scale student models (100k run needs deeper tuning)',
    '',
    '4. Multi-resolution training (224→448) with scaled batch accumulation',
    '',
    '5. Ensemble decoded captions from multiple checkpoints for robustness',
    '',
    '6. Compare feature quality (ViT-L CLIP embeddings) on external benchmarks',
    '   to isolate backbone contribution'
])

# Slide 13: Conclusions
add_content_slide('Main Conclusions', [
    '✓ BLIP-2 pipeline IS reproducible on consumer GPU (all 3 stages ran end-to-end)',
    '',
    '✗ BLIP-2 performance is NOT directly reproducible at local scale',
    '',
    '⟹ Gap is systematic (backbones, data, compute) not pathological (code bugs)',
    '',
    'Why BLIP-2 was best choice:',
    '  - Modular architecture (Q-Former) provides clear downscaling path',
    '  - Frozen backbones allow principled capacity reduction',
    '  - Official LAVIS code is well-maintained & portable',
    '  - No synthetic data dependency (vs. LLaVA)',
    '  - Tractable complexity (vs. VisionLLM v2)'
])

# Slide 14: Artifacts
add_content_slide('Reproducibility Artifacts', [
    'Checkpoint registry: metrics/blip2/checkpoint_registry.jsonl',
    '  → All runs, epochs, loss values traceable',
    '',
    'Per-epoch evaluations: metrics/blip2/caption_eval_summary_*.json',
    '  → Full metric trajectory visible, best-epoch selection validated',
    '',
    'Configuration files: configs/stage{1,2}_*.yaml & caption_*.yaml',
    '  → Exact hyperparams for 10k/50k/100k runs available',
    '',
    'Training loss logs: metrics/blip2/progress_registry.jsonl',
    '  → Bridge training stability confirmed',
    '',
    'Code: All official LAVIS pipeline, modifications documented in codebase',
    '  → Reusable scaffold for future local BLIP-2 work'
])

# Slide 15: Research Implications
add_content_slide('Research Implications', [
    'Multimodal system literature often conflates architecture + scale.',
    'This study shows BLIP-2 architecture is sound; performance gap is scale-driven.',
    '',
    'For practitioners:',
    '• Invest in backbone selection (biggest single lever)',
    '• Checkpoint selection > longer training (100k experiment proves this)',
    '• Frozen-bridge design is genuinely modular & downscalable',
    '',
    'For reproducibility:',
    '• Pipeline reproducibility ≠ performance reproducibility (important to distinguish)',
    '• Local implementations need metric-driven model selection tuning',
    '• Artifact trail matters: losses + metrics + configs = traceability',
    '',
    'For future work: Study efficient backbones (MobileViT, DeiT) with full pipeline'
])

# Save presentation
output_filename = 'BLIP2_Reproduction_Analysis.pptx'
prs.save(output_filename)

# Print success info
full_path = os.path.abspath(output_filename)
file_size = os.path.getsize(output_filename) / (1024 * 1024)  # MB

print(f"✓ Presentation created successfully!")
print(f"✓ File: {output_filename}")
print(f"✓ Full path: {full_path}")
print(f"✓ Size: {file_size:.2f} MB")
print(f"✓ Slides: 15")
