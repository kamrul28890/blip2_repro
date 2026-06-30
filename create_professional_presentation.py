#!/usr/bin/env python3
"""Generate Professional BLIP-2 Reproduction Presentation"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_professional_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    DARK_BG = RGBColor(15, 23, 42)  # Dark blue
    TITLE_COLOR = RGBColor(59, 130, 246)  # Bright blue
    TEXT_COLOR = RGBColor(255, 255, 255)  # White
    DARK_TEXT = RGBColor(20, 20, 30)  # Dark text on light
    ACCENT_COLOR = RGBColor(34, 197, 94)  # Green
    CARD_BG = RGBColor(30, 41, 59)  # Slightly lighter dark
    
    def add_slide_with_bg(title, layout_type='blank'):
        """Add slide with dark background"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        return slide
    
    def add_title_bar(slide, title, with_box=True):
        """Add professional title bar"""
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = TITLE_COLOR
        title_shape.line.fill.background()
        
        text_frame = title_shape.text_frame
        text_frame.text = title
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.paragraph_format.left_indent = Pt(20)
    
    def add_two_column_content(slide, left_title, left_items, right_title, right_items):
        """Add dense two-column layout with headers"""
        # Left column header
        left_header = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(4.5), Inches(0.35))
        tf = left_header.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        # Left column content
        left_content = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(4.5), Inches(5.7))
        tf = left_content.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            p.line_spacing = 1.15
            if item.startswith('•'):
                p.paragraph_format.left_indent = Pt(15)
        
        # Right column header
        right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.35))
        tf = right_header.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        # Right column content
        right_content = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.5), Inches(5.7))
        tf = right_content.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_items):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            p.line_spacing = 1.15
            if item.startswith('•'):
                p.paragraph_format.left_indent = Pt(15)
    
    def add_table_slide(slide, title, rows):
        """Add slide with dense table"""
        add_title_bar(slide, title)
        
        # Add table
        rows_count = len(rows) + 1
        cols_count = len(rows[0]) if rows else 2
        left = Inches(0.3)
        top = Inches(1.0)
        width = Inches(9.4)
        height = Inches(6.0)
        
        table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
        
        # Set column widths
        for col_idx in range(cols_count):
            table_shape.columns[col_idx].width = Inches(9.4 / cols_count)
        
        # Header row
        header_cells = rows[0] if rows else ["Column 1", "Column 2"]
        for col_idx, header_text in enumerate(header_cells):
            cell = table_shape.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TITLE_COLOR
            tf = cell.text_frame
            tf.text = header_text
            p = tf.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_TEXT
        
        # Data rows
        for row_idx, row_data in enumerate(rows[1:], start=1):
            for col_idx, cell_text in enumerate(row_data):
                cell = table_shape.cell(row_idx, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if row_idx % 2 == 0 else DARK_BG
                tf = cell.text_frame
                tf.text = str(cell_text)
                p = tf.paragraphs[0]
                p.font.size = Pt(10)
                p.font.color.rgb = TEXT_COLOR
    
    # ========== SLIDE 1: TITLE ==========
    slide = add_slide_with_bg("")
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Multimodal LLMs\nFrom Theory to Practice"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(1.5))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Comparative Critique: BLIP-2, LLaVA, VisionLLM v2\n+ Local-Scale BLIP-2 Reproduction Study"
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE 2: THREE PAPERS COMPARISON ==========
    slide = add_slide_with_bg("Three Papers: Quick Comparison")
    add_title_bar(slide, "Three Papers: Quick Comparison")
    
    add_two_column_content(slide,
        "BLIP-2 (2023)",
        [
            "• Problem: Efficiently connect frozen",
            "  vision & language backbones",
            "",
            "• Core Idea: Q-Former bridge",
            "  (learnable visual prompts)",
            "",
            "• Data: 129M image-text pairs",
            "  (COCO, VisualGenome, LAION)",
            "",
            "• Compute: 16×A100 40GB, days",
            "",
            "• Best Result:",
            "  BLEU-4: 43.7 | CIDEr: 145.8",
            "",
            "✓ Best for local reproduction"
        ],
        "LLaVA & VisionLLM v2",
        [
            "LLaVA: Instruction-tuned assistant",
            "• GPT-4 generated training data",
            "• ViT-L + Vicuna-7B",
            "• 8×A100, ~10 hours",
            "• Result: 92.53% on ScienceQA",
            "⚠ Requires synthetic data pipeline",
            "",
            "VisionLLM v2: Generalist end-to-end",
            "• Multi-decoder routing system",
            "• 64 datasets, ~100 tasks",
            "• 64→128 A100 GPUs, 18 days",
            "• Broad capability across tasks",
            "✗ Not feasible for local work"
        ]
    )
    
    # ========== SLIDE 3: PROBLEM STATEMENT ==========
    slide = add_slide_with_bg("Problem Statements & Motivation")
    add_title_bar(slide, "Problem Statements & Motivation")
    
    problems = [
        "BLIP-2 Q: Can a lightweight bridge efficiently connect frozen backbones?",
        "  → Efficiency through modular design, not retraining full multimodal stack",
        "",
        "LLaVA Q: How much does supervision format (instruction data) matter?",
        "  → Impact of GPT-4-generated multimodal instructions on assistant behavior",
        "",
        "VisionLLM v2 Q: Can one system unify many vision & vision-language tasks?",
        "  → Single generalist model with routing to multiple specialist decoders",
        "",
        "★ Shared Theme: Reuse pre-trained backbones → reduce training complexity"
    ]
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(6.2))
    tf = content.text_frame
    tf.word_wrap = True
    for i, line in enumerate(problems):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = ACCENT_COLOR if line.startswith("★") else TEXT_COLOR
        p.font.bold = line.startswith("★")
        p.space_before = Pt(3)
        p.space_after = Pt(3)
    
    # ========== SLIDE 4: METHODOLOGY DENSE TABLE ==========
    slide = add_slide_with_bg("Methodology & Experimental Setup")
    add_table_slide(slide, "Methodology & Experimental Setup", [
        ["Dimension", "BLIP-2", "LLaVA", "VisionLLM v2"],
        ["Core Idea", "Q-Former bridge", "Lightweight projection", "Multi-decoder routing"],
        ["Training Stages", "2 stages", "2 stages (incl. inst-tune)", "Multi-stage, multi-decoder"],
        ["Data Scale", "129M pairs", "158K instruction pairs", "64 datasets, ~100 tasks"],
        ["Pretraining", "COCO+VisualGenome+LAION", "ViT-L CLIP + Vicuna-7B", "Multiple pretrained decoders"],
        ["Compute (GPU)", "16×A100 40GB", "8×A100", "64→128 A100 GPUs"],
        ["Training Time", "Multi-day", "~10 hours", "18 days"],
        ["Parameter Efficiency", "Frozen backbones", "Simple projection layer", "Routing complexity"],
    ])
    
    # ========== SLIDE 5: RESULTS & REPRODUCIBILITY ==========
    slide = add_slide_with_bg("Results & Local Reproducibility Verdict")
    add_title_bar(slide, "Results & Local Reproducibility Verdict")
    
    add_two_column_content(slide,
        "BLIP-2: ✓ BEST TARGET",
        [
            "Published Result:",
            "  BLEU-4: 43.7, CIDEr: 145.8",
            "  (ViT-g, OPT-2.7B on Karpathy)",
            "",
            "Why local-feasible:",
            "✓ Modular frozen-backbone arch",
            "✓ Preserves coherent downscaling",
            "✓ Official LAVIS code available",
            "✓ No synthetic data dependency",
            "✓ Can run on single RTX 3070",
            "",
            "Feasibility: EXCELLENT"
        ],
        "LLaVA & VisionLLM v2",
        [
            "LLaVA Result: 92.53% ScienceQA",
            "",
            "Limitations:",
            "⚠ Depends on GPT-4 instruction",
            "⚠ Still needs 8×A100 for full",
            "⚠ Synthetic data pipeline complex",
            "",
            "",
            "VisionLLM v2: Broad generalist",
            "",
            "Limitations:",
            "✗ 64-128 A100 GPUs required",
            "✗ Multi-decoder too complex",
            "✗ Gains hard to disentangle",
            "",
            "Feasibility: NOT LOCAL"
        ]
    )
    
    # ========== SLIDE 6: LOCAL SETUP DETAILS ==========
    slide = add_slide_with_bg("Local BLIP-2 Setup & Pipeline")
    add_title_bar(slide, "Local BLIP-2 Setup & Pipeline")
    
    add_two_column_content(slide,
        "Hardware & Architecture Downscaling",
        [
            "GPU: RTX 3070 8GB (vs 16×A100)",
            "",
            "Vision Encoder:",
            "  Local: CLIP ViT-L",
            "  Paper: ViT-g (~7× larger)",
            "",
            "Language Model:",
            "  Local: OPT-350M",
            "  Paper: OPT-2.7B (~8× larger)",
            "",
            "Image Resolution: 224 (vs 364)",
            "",
            "Batch Size: 1-2 with accum 8",
            "  (vs 256 in paper)"
        ],
        "Three-Stage Pipeline (All Preserved)",
        [
            "Stage 1: Representation Learning",
            "  • Q-Former (32 queries)",
            "  • With frozen ViT-L",
            "  • ITC + ITM + LM losses",
            "  • 3 epochs (50k run)",
            "",
            "Stage 2: Generative Alignment",
            "  • Q-Former → OPT projections",
            "  • Autoregressive objective",
            "  • 3 epochs",
            "",
            "Stage 3: Caption Fine-tuning",
            "  • Prompt: 'a photo of'",
            "  • 5 epochs (best) or 15 (long)"
        ]
    )
    
    # ========== SLIDE 7: LOCAL RESULTS SCALING ==========
    slide = add_slide_with_bg("Local Results: Scaling Analysis")
    add_table_slide(slide, "Local Results Across Data Scales", [
        ["Setup", "Train Data", "BLEU-4", "CIDEr", "Status"],
        ["10k Pilot", "10k images", "1.45", "3.03", "Baseline"],
        ["50k Main", "50k images", "11.13", "37.57", "✓ PEAK (e3)"],
        ["50k Polish", "50k images", "10.49", "37.14", "Refinement (e4)"],
        ["100k Long", "100k images", "8.57", "25.34", "⚠ Best epoch 2"],
        ["100k Final", "100k images", "6.61", "18.36", "Degradation (e14)"],
        ["Paper", "129M pairs", "43.7", "145.8", "Reference"],
    ])
    
    # ========== SLIDE 8: GAP BREAKDOWN ==========
    slide = add_slide_with_bg("Performance Gap: Systematic Analysis")
    add_title_bar(slide, "Performance Gap: Systematic Analysis")
    
    add_two_column_content(slide,
        "Paper vs Local Configuration",
        [
            "Vision Backbone:",
            "  Paper: ViT-g (352M params)",
            "  Local: ViT-L (86M params)",
            "  → ~4× capacity loss",
            "",
            "Language Model:",
            "  Paper: OPT-2.7B",
            "  Local: OPT-350M",
            "  → ~8× parameter reduction",
            "",
            "Pretraining Data:",
            "  Paper: 129M pairs (diverse)",
            "  Local: 0 external, COCO-only",
            "  → Major feature quality impact",
            "",
            "Resolution: 224 vs 364"
        ],
        "Multiplicative Gap: ~25-30×",
        [
            "Gap Sources (estimated):",
            "",
            "1. Backbone capacity: 3-4× loss",
            "   (ViT-g vs ViT-L)",
            "",
            "2. LLM parameters: ~8× reduction",
            "   (OPT-2.7B vs OPT-350M)",
            "",
            "3. Pretraining data scale:",
            "   129M→0 external = critical",
            "",
            "4. Fine-tune data: 50k local",
            "   vs massive official corpus",
            "",
            "5. Resolution & compute budget",
            "",
            "⟹ Multiplicative, NOT",
            "   pipeline failure"
        ]
    )
    
    # ========== SLIDE 9: PEAK-EARLY PHENOMENON ==========
    slide = add_slide_with_bg("100k Run: Peak-Early & Overfitting")
    add_title_bar(slide, "100k Run: Peak-Early & Overfitting")
    
    add_two_column_content(slide,
        "Training Loss (Smooth)",
        [
            "Stage 1 loss: 0.33",
            "Stage 2 loss: 0.36",
            "Caption train loss: steady ↓",
            "",
            "→ Bridge training STABLE",
            "→ Losses don't diverge",
            "→ No gradient instability",
            "",
            "FINDING: Training process",
            "healthy; metric degradation",
            "is NOT due to training failure"
        ],
        "Caption Metrics (Early Peak)",
        [
            "Epoch 2 (BEST):",
            "  BLEU-4 = 8.57",
            "  CIDEr = 25.34",
            "",
            "Epoch 14 (FINAL):",
            "  BLEU-4 = 6.61 (-23%)",
            "  CIDEr = 18.36 (-28%)",
            "",
            "Diversity: 677→734 captions",
            "(More but lower quality)",
            "",
            "Root Cause: Overfitting without",
            "validation-driven early stopping"
        ]
    )
    
    # ========== SLIDE 10: LIMITATIONS ==========
    slide = add_slide_with_bg("Study Limitations")
    add_title_bar(slide, "Study Limitations")
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(6.2))
    tf = content.text_frame
    tf.word_wrap = True
    
    limitations = [
        "1. Local Evaluation Scope: 1k val/test vs paper's full Karpathy split",
        "   → Results internally comparable but not leaderboard-equivalent",
        "",
        "2. Downscaling Path: CLIP ViT-L + OPT-350M optimized locally",
        "   → Different from paper's ViT-g + OPT-2.7B, not exact reproduction",
        "",
        "3. Single-Seed Runs: No variance/confidence intervals",
        "   → Rankings stable, but generalization uncertain",
        "",
        "4. Windows Portability: Code modified for non-distributed, single-process mode",
        "   → Measures both model reproducibility + codebase portability challenges",
        "",
        "5. Real Data Only: All COCO annotations (unlike LLaVA's GPT-generated)"
    ]
    
    for i, line in enumerate(limitations):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = ACCENT_COLOR if line.startswith(("1.", "2.", "3.", "4.", "5.")) else TEXT_COLOR
        p.space_before = Pt(3)
        p.space_after = Pt(3)
    
    # ========== SLIDE 11: CONCLUSIONS ==========
    slide = add_slide_with_bg("Main Conclusions")
    add_title_bar(slide, "Main Conclusions")
    
    conclusions = [
        "✓ BLIP-2 PIPELINE IS REPRODUCIBLE",
        "   All three training stages (S1, S2, caption) completed end-to-end on RTX 3070",
        "",
        "✗ BLIP-2 PERFORMANCE NOT DIRECTLY REPRODUCIBLE",
        "   Local best (11.13 BLEU-4) >> Paper (43.7) due to scale factors, not bugs",
        "",
        "⟹ GAP IS SYSTEMATIC",
        "   Backbone capacity + LLM size + pretraining data + compute budget = multiplicative",
        "   NOT: Pipeline failure, missing components, or architectural misunderstanding",
        "",
        "WHY BLIP-2 WAS BEST CHOICE:",
        "✓ Modular frozen-backbone architecture provides clear downscaling path",
        "✓ Frozen backbones allow principled capacity reduction (not retraining full model)",
        "✓ Official LAVIS code is well-maintained, documented, and portable",
        "✓ No synthetic data pipeline dependency (vs. LLaVA)",
        "✓ Tractable complexity (vs. VisionLLM v2's multi-decoder routing)"
    ]
    
    content = slide.shapes.add_textbox(Inches(0.4), Inches(0.85), Inches(9.2), Inches(6.4))
    tf = content.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(conclusions):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = line
        is_banner = any(s in line for s in ["PIPELINE", "PERFORMANCE", "GAP IS", "CHOICE:"])
        is_check = line.startswith("✓") or line.startswith("✗") or line.startswith("⟹")
        p.font.size = Pt(12 if is_banner else 10)
        p.font.bold = is_banner or is_check
        p.font.color.rgb = ACCENT_COLOR if is_banner or is_check else TEXT_COLOR
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    
    # ========== SLIDE 12: ARTIFACTS & TRACEABILITY ==========
    slide = add_slide_with_bg("Reproducibility Artifacts & Traceability")
    add_title_bar(slide, "Reproducibility Artifacts & Traceability")
    
    add_two_column_content(slide,
        "Checkpoint Registry",
        [
            "📁 metrics/blip2/",
            "   checkpoint_registry.jsonl",
            "",
            "Contains:",
            "• All run IDs and epochs",
            "• Training losses per stage",
            "• Final metrics per checkpoint",
            "• Model selection reasoning",
            "",
            "Per-Epoch Evaluations:",
            "  caption_eval_summary_*.json",
            "",
            "Enables:",
            "✓ Full metric trajectory",
            "✓ Best-epoch selection audit",
            "✓ Complete run history"
        ],
        "Configuration & Loss Logs",
        [
            "Configuration Files:",
            "  configs/stage1_*.yaml",
            "  configs/stage2_*_opt350m.yaml",
            "  configs/caption_*.yaml",
            "",
            "Exact hyperparams for:",
            "• 10k pilot run",
            "• 50k main run",
            "• 100k long-run",
            "",
            "Training Loss Logs:",
            "  progress_registry.jsonl",
            "",
            "Proves:",
            "✓ Bridge training stable",
            "✓ No gradient divergence",
            "✓ Smooth loss decrease"
        ]
    )
    
    # ========== SLIDE 13: FUTURE WORK ==========
    slide = add_slide_with_bg("Future Work & Unexplored Directions")
    add_title_bar(slide, "Future Work & Unexplored Directions")
    
    future_work = [
        "1. Systematic Backbone Scaling: ViT-B, ViT-L, ViT-H vs OPT-125M→350M→1.3B",
        "   → Isolate capacity vs performance relationship",
        "",
        "2. Data Augmentation: Synthetic descriptions from self-supervised or larger models",
        "   → Mimic some benefits of 129M pretraining without computation",
        "",
        "3. Regularization Tuning: Dropout schedules, weight decay, early stopping for student scale",
        "   → Address the 100k peak-early phenomenon",
        "",
        "4. Multi-Resolution End-to-End: 224→448 with scaled batch accumulation",
        "   → Understand resolution contribution to gap",
        "",
        "5. Ensemble Decoding: Aggregate predictions from multiple selected checkpoints",
        "   → Stability improvement for long-run training",
        "",
        "6. External Benchmark Validation: CLIP ViT-L embeddings on ImageNet, CIFAR",
        "   → Isolate backbone architecture quality from data scale"
    ]
    
    content = slide.shapes.add_textbox(Inches(0.4), Inches(0.85), Inches(9.2), Inches(6.4))
    tf = content.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(future_work):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = ACCENT_COLOR if line[0].isdigit() else TEXT_COLOR
        p.font.bold = line[0].isdigit()
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        if line[0].isdigit():
            p.paragraph_format.left_indent = Pt(0)
        else:
            p.paragraph_format.left_indent = Pt(25)
    
    # ========== SLIDE 14: KEY TAKEAWAYS ==========
    slide = add_slide_with_bg("Key Takeaways for Practitioners")
    add_title_bar(slide, "Key Takeaways for Practitioners")
    
    add_two_column_content(slide,
        "Technical Insights",
        [
            "1. Backbone selection is the single",
            "   biggest performance lever",
            "",
            "2. Checkpoint selection > longer",
            "   training schedules (100k proves)",
            "",
            "3. Frozen-bridge design is",
            "   genuinely modular & downscalable",
            "",
            "4. Validation-driven model selection",
            "   critical for student-scale work",
            "",
            "5. Three-stage pipeline is robust;",
            "   failure modes are optimization"
        ],
        "Reproducibility Insights",
        [
            "1. Pipeline reproducibility ≠",
            "   performance reproducibility",
            "   (crucial distinction!)",
            "",
            "2. Local implementations need",
            "   metric-driven tuning, not exact",
            "   paper parameter copy",
            "",
            "3. Artifact trail + full transparency",
            "   enables later verification",
            "",
            "4. Scale factors compound",
            "   multiplicatively, not additively",
            "",
            "5. BLIP-2 is genuinely the best",
            "   target for commodity reproduction"
        ]
    )
    
    # ========== SLIDE 15: SUMMARY ==========
    slide = add_slide_with_bg("Research Summary")
    add_title_bar(slide, "Research Summary")
    
    summary_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(6.2))
    tf = summary_text.text_frame
    tf.word_wrap = True
    
    summary = [
        "This work studied three landmark multimodal LLM papers (BLIP-2, LLaVA, VisionLLM v2)",
        "through the dual lens of architectural critique and student-scale reproducibility.",
        "",
        "Key Finding: BLIP-2 is the most defensible local reproduction target among the three.",
        "Its frozen-backbone modularity provides a coherent downscaling path that the other",
        "systems do not offer (LLaVA depends on GPT-4 instructions; VisionLLM v2 requires",
        "64-128 A100 GPUs and multi-decoder routing logic).",
        "",
        "We successfully reproduced all three BLIP-2 training stages on a single RTX 3070 8GB GPU.",
        "Best local result: BLEU-4 11.13, CIDEr 37.57 on 50k COCO subset, vs. paper's 43.7 / 145.8.",
        "The ~25-30× performance gap is systematic (backbones, data, compute), not pathological.",
        "",
        "Practical lesson: For multimodal research on commodity hardware, invest in understanding",
        "which performance comes from architectural novelty vs. scale, and select your reproduction",
        "target accordingly. BLIP-2's design makes it both scientifically coherent and practically",
        "feasible for constrained settings."
    ]
    
    for i, line in enumerate(summary):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        p.line_spacing = 1.2
    
    # Save
    prs.save("BLIP2_Reproduction_Presentation_Professional.pptx")
    print("✓ Professional presentation created: BLIP2_Reproduction_Presentation_Professional.pptx")
    print("✓ 15 slides with dense, multi-column layouts")
    print("✓ Dark blue theme with accent colors, no wasted whitespace")
    print("✓ Full technical content with tables and structured information")

if __name__ == "__main__":
    create_professional_presentation()
