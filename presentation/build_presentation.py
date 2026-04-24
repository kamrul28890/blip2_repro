from __future__ import annotations

import json
from pathlib import Path
from typing import Sequence

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = ROOT / "presentation"
OUTPUT_PATH = PRESENTATION_DIR / "mkamrul_blip2_term_presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TITLE_FONT = "Bahnschrift"
BODY_FONT = "Calibri"


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.strip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


BG = rgb("F6F3ED")
HEADER = rgb("224A59")
TEXT = rgb("1C2830")
MUTED = rgb("58646E")
ACCENT = rgb("C9782B")
GREEN = rgb("2E6B57")
RED = rgb("AA4E45")
BLUE = rgb("6A8BA3")
PALE_BLUE = rgb("DCEAF1")
PALE_GREEN = rgb("DDEEE6")
PALE_RED = rgb("F3E0DD")
PALE_GOLD = rgb("F5E7D4")
CARD = rgb("FEFEFD")
LINE = rgb("C9D2D6")


def load_json(relative_path: str) -> dict:
    return json.loads((ROOT / relative_path).read_text(encoding="utf-8"))


REPORT = load_json("metrics/blip2/report_metrics_snapshot.json")
RUN_100K_E2 = load_json("metrics/blip2/caption_eval_summary_student_100k_long_epoch2.json")
RUN_100K_E14 = load_json("metrics/blip2/caption_eval_summary_student_100k_long_epoch14.json")
QUAL_EXAMPLES = load_json("metrics/blip2/caption_eval_examples_office_epoch3.json")


def set_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.72))
    band.fill.solid()
    band.fill.fore_color.rgb = HEADER
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.13), Inches(9.9), Inches(0.3))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = TITLE_FONT
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = rgb("FFFFFF")

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(9.7), Inches(0.14), Inches(3.1), Inches(0.24))
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        p.alignment = PP_ALIGN.RIGHT
        p.font.name = BODY_FONT
        p.font.size = Pt(11)
        p.font.color.rgb = rgb("D8E3E7")


def add_footer(slide, source_text: str, slide_no: int) -> None:
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.32), Inches(7.02), Inches(12.65), Inches(0.01)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.06), Inches(11.5), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = source_text
    p.font.name = BODY_FONT
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED

    num = slide.shapes.add_textbox(Inches(12.15), Inches(7.02), Inches(0.65), Inches(0.2))
    p = num.text_frame.paragraphs[0]
    p.text = str(slide_no)
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = BODY_FONT
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = HEADER


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: Sequence[str],
    font_size: int = 16,
    color: RGBColor = TEXT,
    bold_first: bool = False,
    bullet: bool = False,
    line_spacing: float = 1.05,
    left: float = 0,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for idx, text in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = BODY_FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold_first and idx == 0
        p.space_after = Pt(3)
        p.line_spacing = line_spacing
        p.level = 0
        p.bullet = bullet
        if left:
            p.left_margin = Pt(left)
    return box


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body_lines: Sequence[str],
    fill_rgb: RGBColor = CARD,
    title_rgb: RGBColor = HEADER,
    body_rgb: RGBColor = TEXT,
    font_size: int = 15,
    title_size: int = 18,
):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.1)

    title_box = slide.shapes.add_textbox(x + Inches(0.16), y + Inches(0.1), w - Inches(0.32), Inches(0.42))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = TITLE_FONT
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = title_rgb

    add_textbox(
        slide,
        x + Inches(0.14),
        y + Inches(0.52),
        w - Inches(0.28),
        h - Inches(0.62),
        body_lines,
        font_size=font_size,
        color=body_rgb,
        bullet=True,
        left=14,
    )
    return shape


def add_label(slide, x: float, y: float, w: float, h: float, text: str, fill_rgb: RGBColor, text_rgb: RGBColor):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = BODY_FONT
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = text_rgb
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow(slide, x: float, y: float, w: float, h: float, fill_rgb: RGBColor):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def add_metric_box(slide, x: float, y: float, w: float, h: float, label: str, value: str, fill_rgb: RGBColor):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED
    p = tf.add_paragraph()
    p.text = value
    p.alignment = PP_ALIGN.CENTER
    p.font.name = TITLE_FONT
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = HEADER
    return shape


def add_table(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    rows: Sequence[Sequence[str]],
    col_widths: Sequence[float] | None = None,
    header_fill: RGBColor = HEADER,
    header_text: RGBColor = rgb("FFFFFF"),
    row_fills: Sequence[RGBColor] | None = None,
    font_size: int = 12,
):
    table = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h).table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = width

    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.font.name = BODY_FONT
            p.font.size = Pt(font_size)
            p.font.color.rgb = TEXT
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                p.font.color.rgb = header_text
                p.font.bold = True
            else:
                fill = row_fills[(r - 1) % len(row_fills)] if row_fills else CARD
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill
    return table


def add_chart(
    slide,
    chart_type,
    x: float,
    y: float,
    w: float,
    h: float,
    categories: Sequence[str],
    series: Sequence[tuple[str, Sequence[float], RGBColor]],
    y_max: float | None = None,
    legend: bool = True,
):
    chart_data = CategoryChartData()
    chart_data.categories = list(categories)
    for name, values, _color in series:
        chart_data.add_series(name, list(values))

    chart = slide.shapes.add_chart(chart_type, x, y, w, h, chart_data).chart
    chart.has_title = False
    chart.has_legend = legend
    if legend and chart.legend is not None:
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.name = BODY_FONT
        chart.legend.font.size = Pt(10)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = LINE
    chart.value_axis.tick_labels.font.name = BODY_FONT
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.name = BODY_FONT
    chart.category_axis.tick_labels.font.size = Pt(10)
    if y_max is not None:
        chart.value_axis.maximum_scale = y_max
    chart.value_axis.minimum_scale = 0
    for idx, (_, _values, color) in enumerate(series):
        chart.series[idx].format.fill.solid()
        chart.series[idx].format.fill.fore_color.rgb = color
        chart.series[idx].format.line.color.rgb = color
    return chart


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    img = Image.open(path)
    iw, ih = img.size
    box_ratio = float(w) / float(h)
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        new_w = w
        new_h = int(float(w) / img_ratio)
        top = y + int((float(h) - new_h) / 2)
        left = x
    else:
        new_h = h
        new_w = int(float(h) * img_ratio)
        left = x + int((float(w) - new_w) / 2)
        top = y
    return slide.shapes.add_picture(str(path), left, top, width=new_w, height=new_h)


def add_progression_box(slide, x: float, y: float, w: float, h: float, title: str, body: str, fill_rgb: RGBColor):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.name = TITLE_FONT
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = rgb("FFFFFF")
    p = tf.add_paragraph()
    p.text = body
    p.alignment = PP_ALIGN.CENTER
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.color.rgb = rgb("F4F8FA")
    return shape


def new_slide(prs: Presentation, title: str, slide_no: int, source_text: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, title, subtitle=subtitle)
    add_footer(slide, source_text, slide_no)
    return slide


def build_presentation() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_no = 1

    slide = new_slide(
        prs,
        "Three Papers, One Multimodal Progression",
        slide_no,
        "Sources: BLIP-2 (2023), LLaVA (2023), VisionLLM v2 (2024)",
    )
    add_label(slide, Inches(0.72), Inches(1.0), Inches(2.05), Inches(0.28), "2023", PALE_BLUE, HEADER)
    add_card(
        slide,
        Inches(0.72),
        Inches(1.25),
        Inches(3.8),
        Inches(2.65),
        "BLIP-2",
        [
            "Freezes the image encoder and LLM, then trains a lightweight Q-Former bridge.",
            "Main goal: make strong frozen backbones talk to each other efficiently.",
            "Best fit when we care about architecture and staged learning.",
        ],
        fill_rgb=CARD,
    )
    add_label(slide, Inches(4.78), Inches(1.0), Inches(2.05), Inches(0.28), "2023", PALE_GOLD, ACCENT)
    add_card(
        slide,
        Inches(4.78),
        Inches(1.25),
        Inches(3.8),
        Inches(2.65),
        "LLaVA",
        [
            "Connects vision features to a Vicuna-like LLM, then instruction-tunes on GPT-generated visual dialogs.",
            "Main goal: multimodal assistant behavior and chat-style interaction.",
            "Best fit when we care about usability and instruction following.",
        ],
        fill_rgb=CARD,
    )
    add_label(slide, Inches(8.84), Inches(1.0), Inches(2.05), Inches(0.28), "2024", PALE_GREEN, GREEN)
    add_card(
        slide,
        Inches(8.84),
        Inches(1.25),
        Inches(3.78),
        Inches(2.65),
        "VisionLLM v2",
        [
            "Uses an MLLM plus a super-link mechanism to drive multiple task-specific decoders.",
            "Main goal: one system for perception, understanding, generation, and editing.",
            "Best fit when we care about breadth across many vision-language tasks.",
        ],
        fill_rgb=CARD,
    )
    add_arrow(slide, Inches(4.28), Inches(2.2), Inches(0.42), Inches(0.35), BLUE)
    add_arrow(slide, Inches(8.34), Inches(2.2), Inches(0.42), Inches(0.35), ACCENT)
    add_card(
        slide,
        Inches(0.72),
        Inches(4.25),
        Inches(11.92),
        Inches(1.85),
        "Common Thread",
        [
            "All three papers start from the same big question: how do we turn strong vision and language backbones into practical multimodal systems?",
            "The progression is bridge learning -> instruction-following assistants -> multi-output generalist systems.",
            "As the systems become broader and more capable, they also become harder to reproduce on student-scale hardware.",
        ],
        fill_rgb=PALE_BLUE,
    )
    slide_no += 1

    slide = new_slide(
        prs,
        "How The Papers Connect",
        slide_no,
        "Source lens: paper critique in paper/term_paper.tex",
    )
    rows = [
        ["Question", "BLIP-2", "LLaVA", "VisionLLM v2"],
        ["Where alignment happens", "Q-Former bridge", "Projection + instruction tuning", "MLLM backbone + super link"],
        ["Main supervision", "Image-text pretraining", "GPT-generated multimodal instructions", "Hundreds of task datasets"],
        ["Typical output", "Text answers / captions", "Chat-style text responses", "Text + boxes + pose + generated images"],
        ["Best at", "Modular transfer from frozen backbones", "Assistant behavior and interactive reasoning", "Broad cross-task generality"],
        ["Reproducibility on class hardware", "Best of the three", "Possible but still costly", "Least realistic"],
    ]
    add_table(
        slide,
        Inches(0.6),
        Inches(1.02),
        Inches(12.1),
        Inches(4.0),
        rows,
        col_widths=[Inches(2.55), Inches(3.1), Inches(3.05), Inches(3.4)],
        row_fills=[CARD, PALE_BLUE],
        font_size=12,
    )
    add_card(
        slide,
        Inches(0.85),
        Inches(5.32),
        Inches(11.45),
        Inches(1.28),
        "Connection In One Sentence",
        [
            "BLIP-2 focuses on a clean bridge, LLaVA focuses on instruction-following behavior, and VisionLLM v2 focuses on widening the output space of multimodal models.",
            "That makes them easy to compare in proper order: connection mechanism first, training signal second, and task scope third.",
        ],
        fill_rgb=PALE_GOLD,
    )
    slide_no += 1

    slide = new_slide(
        prs,
        "Why BLIP-2 Was The Right Paper To Reproduce",
        slide_no,
        "Sources: paper/mkamrul_project_answers.md, paper/term_paper.tex",
    )
    rows = [
        ["Criterion", "BLIP-2", "LLaVA", "VisionLLM v2"],
        ["Architectural clarity", "Strong", "Medium", "Medium"],
        ["Public implementation", "Strong", "Strong", "Medium"],
        ["Downscales meaningfully", "Strong", "Medium", "Weak"],
        ["Single-GPU educational value", "Strong", "Medium", "Weak"],
        ["Risk of project collapse", "Lowest", "Moderate", "Highest"],
    ]
    add_table(
        slide,
        Inches(0.65),
        Inches(1.15),
        Inches(7.05),
        Inches(4.75),
        rows,
        col_widths=[Inches(2.15), Inches(1.55), Inches(1.55), Inches(1.8)],
        row_fills=[CARD, PALE_BLUE],
        font_size=13,
    )
    add_card(
        slide,
        Inches(8.0),
        Inches(1.18),
        Inches(4.65),
        Inches(2.18),
        "Selection Logic",
        [
            "BLIP-2 still makes sense after aggressive scaling because the trainable part is the bridge, not the entire model.",
            "The official LAVIS code path exists and can be audited, patched, and reused.",
            "It stays non-trivial: three stages, custom evaluation, and real systems work were still required.",
        ],
        fill_rgb=PALE_GREEN,
        title_rgb=GREEN,
    )
    add_card(
        slide,
        Inches(8.0),
        Inches(3.62),
        Inches(4.65),
        Inches(1.95),
        "Project Thesis",
        [
            "Target: methodological fidelity under hardware limits, not a literal paper-scale replication.",
            "Question: can we reproduce the pipeline even if we cannot reproduce the published score?",
        ],
        fill_rgb=PALE_RED,
        title_rgb=RED,
    )
    slide_no += 1

    slide = new_slide(prs, "BLIP-2 In A Nutshell", slide_no, "Sources: BLIP-2 paper, term_paper.tex")
    add_label(slide, Inches(0.85), Inches(1.02), Inches(2.05), Inches(0.28), "Frozen", PALE_BLUE, HEADER)
    add_card(slide, Inches(0.75), Inches(1.3), Inches(2.0), Inches(1.28), "Image", ["Raw image enters a pretrained vision encoder."], font_size=14)
    add_card(slide, Inches(2.95), Inches(1.3), Inches(2.0), Inches(1.28), "Image Encoder", ["A frozen visual backbone produces image features."], fill_rgb=PALE_BLUE, font_size=14)
    add_label(slide, Inches(5.18), Inches(1.02), Inches(2.05), Inches(0.28), "Trainable", PALE_GOLD, ACCENT)
    add_card(slide, Inches(5.1), Inches(1.3), Inches(2.05), Inches(1.28), "Q-Former", ["32 learned queries pull language-relevant visual information."], fill_rgb=PALE_GOLD, title_rgb=ACCENT, font_size=14)
    add_card(slide, Inches(7.42), Inches(1.3), Inches(1.65), Inches(1.28), "Projection", ["Maps query outputs into the language model embedding space."], fill_rgb=CARD, font_size=13)
    add_card(slide, Inches(9.3), Inches(1.3), Inches(2.05), Inches(1.28), "Frozen LLM", ["The language model stays frozen and generates text from bridged inputs."], fill_rgb=PALE_BLUE, font_size=14)
    add_card(slide, Inches(11.58), Inches(1.3), Inches(1.0), Inches(1.28), "Output", ["Caption or answer"], fill_rgb=CARD, font_size=13)
    add_arrow(slide, Inches(2.52), Inches(1.73), Inches(0.28), Inches(0.22), HEADER)
    add_arrow(slide, Inches(4.72), Inches(1.73), Inches(0.28), Inches(0.22), ACCENT)
    add_arrow(slide, Inches(7.1), Inches(1.73), Inches(0.28), Inches(0.22), HEADER)
    add_arrow(slide, Inches(9.02), Inches(1.73), Inches(0.28), Inches(0.22), HEADER)
    add_arrow(slide, Inches(11.28), Inches(1.73), Inches(0.22), Inches(0.22), HEADER)
    add_card(
        slide,
        Inches(0.78),
        Inches(3.1),
        Inches(5.7),
        Inches(2.58),
        "Close Domain Work",
        [
            "BLIP (2022): unified vision-language pretraining for understanding and generation.",
            "Flamingo (2022): large multimodal LM that interleaves visual tokens with language modeling at scale.",
            "BLIP-2 sits between them: it keeps powerful backbones frozen and concentrates learning in a compact bridge.",
        ],
        fill_rgb=CARD,
    )
    add_card(
        slide,
        Inches(6.82),
        Inches(3.1),
        Inches(5.8),
        Inches(2.58),
        "What Makes BLIP-2 Different",
        [
            "It makes the cross-modal interface explicit through Q-Former rather than fully tuning the whole stack.",
            "That gives it a cleaner experimental story: if results move, the bridge is the main thing changing.",
            "This is exactly why it became the best reproduction target for the project.",
        ],
        fill_rgb=PALE_BLUE,
    )
    slide_no += 1

    slide = new_slide(prs, "BLIP-2 Method, Strengths, Limits", slide_no, "Sources: BLIP-2 paper, term_paper.tex")
    add_card(
        slide,
        Inches(0.7),
        Inches(1.05),
        Inches(3.9),
        Inches(1.58),
        "Stage 1: Representation Learning",
        [
            "Learns language-relevant visual queries from paired image-text data.",
            "Uses image-text contrastive, matching, and language-modeling style losses.",
        ],
        fill_rgb=PALE_BLUE,
    )
    add_card(
        slide,
        Inches(4.72),
        Inches(1.05),
        Inches(3.9),
        Inches(1.58),
        "Stage 2: Generative Alignment",
        [
            "Aligns the learned query outputs to the frozen LLM for text generation.",
            "Turns the bridge from a representation learner into a generation interface.",
        ],
        fill_rgb=PALE_GOLD,
        title_rgb=ACCENT,
    )
    add_card(
        slide,
        Inches(8.74),
        Inches(1.05),
        Inches(3.9),
        Inches(1.58),
        "Task Fine-Tuning",
        [
            "Adapts the aligned model to captioning, VQA, retrieval, and other downstream tasks.",
            "This is where the paper reports strong captioning and zero-shot behavior.",
        ],
        fill_rgb=PALE_GREEN,
        title_rgb=GREEN,
    )
    add_card(
        slide,
        Inches(0.85),
        Inches(3.05),
        Inches(5.65),
        Inches(2.42),
        "Where BLIP-2 Is Strong",
        [
            "Strong conceptual modularity: frozen backbones + small bridge.",
            "Good transfer story across captioning, retrieval, and VQA.",
            "Much easier to reason about than a fully coupled multimodal stack.",
        ],
        fill_rgb=CARD,
    )
    add_card(
        slide,
        Inches(6.84),
        Inches(3.05),
        Inches(5.65),
        Inches(2.42),
        "Main Limitations",
        [
            "The paper's 'efficient' setup is still large: 129M pretraining images and large GPU budgets.",
            "Best published results use much bigger backbones than a student desktop can fit.",
            "So the idea is downscalable, but the top performance is still scale-sensitive.",
        ],
        fill_rgb=PALE_RED,
        title_rgb=RED,
    )
    slide_no += 1

    slide = new_slide(prs, "LLaVA In A Nutshell", slide_no, "Sources: LLaVA paper, term_paper.tex")
    add_label(slide, Inches(0.8), Inches(1.02), Inches(2.4), Inches(0.28), "Instruction-Tuned Assistant", PALE_GOLD, ACCENT)
    add_card(slide, Inches(0.8), Inches(1.32), Inches(1.8), Inches(1.18), "Image", ["User image"], font_size=14)
    add_card(slide, Inches(2.95), Inches(1.32), Inches(2.15), Inches(1.18), "Vision Encoder", ["Turns the image into visual features."], fill_rgb=PALE_BLUE, font_size=14)
    add_card(slide, Inches(5.42), Inches(1.32), Inches(1.55), Inches(1.18), "Projector", ["Maps vision features into the LLM input space."], fill_rgb=PALE_GOLD, title_rgb=ACCENT, font_size=13)
    add_card(slide, Inches(7.28), Inches(1.32), Inches(2.1), Inches(1.18), "Vicuna-like LLM", ["Generates open-ended chat responses."], fill_rgb=PALE_BLUE, font_size=14)
    add_card(slide, Inches(9.7), Inches(1.32), Inches(2.25), Inches(1.18), "Instruction Data", ["GPT-4 is used to synthesize multimodal training conversations."], fill_rgb=PALE_GREEN, title_rgb=GREEN, font_size=13)
    add_card(slide, Inches(10.25), Inches(2.88), Inches(2.1), Inches(1.0), "Output", ["Chat answer, explanation, reasoning"], font_size=13)
    add_arrow(slide, Inches(2.58), Inches(1.67), Inches(0.28), Inches(0.22), HEADER)
    add_arrow(slide, Inches(5.06), Inches(1.67), Inches(0.28), Inches(0.22), ACCENT)
    add_arrow(slide, Inches(6.94), Inches(1.67), Inches(0.28), Inches(0.22), HEADER)
    add_card(
        slide,
        Inches(0.82),
        Inches(4.05),
        Inches(5.75),
        Inches(1.85),
        "What It Is Really Trying To Do",
        [
            "LLaVA is less about a new multi-stage pretraining recipe and more about turning a vision-language model into a usable assistant.",
            "Its key training move is visual instruction tuning, not just image-text alignment.",
        ],
        fill_rgb=CARD,
    )
    add_card(
        slide,
        Inches(6.82),
        Inches(4.05),
        Inches(5.7),
        Inches(1.85),
        "How It Differs From BLIP-2",
        [
            "BLIP-2 asks whether a small bridge can unlock frozen backbones.",
            "LLaVA asks whether instruction data can make a multimodal model behave like a chat assistant.",
        ],
        fill_rgb=PALE_BLUE,
    )
    slide_no += 1

    slide = new_slide(prs, "LLaVA Method, Strengths, Limits", slide_no, "Sources: LLaVA paper, term_paper.tex")
    add_card(
        slide,
        Inches(0.7),
        Inches(1.05),
        Inches(3.88),
        Inches(2.35),
        "Method",
        [
            "Stage 1 aligns visual features to the LLM.",
            "Stage 2 instruction-tunes the system on GPT-generated multimodal dialog data.",
            "Benchmarking emphasizes chat quality and instruction following.",
        ],
        fill_rgb=PALE_GOLD,
        title_rgb=ACCENT,
    )
    add_card(
        slide,
        Inches(4.74),
        Inches(1.05),
        Inches(3.88),
        Inches(2.35),
        "Where It Is Good",
        [
            "Explaining images in natural language.",
            "Interactive question answering and assistant-like responses.",
            "Strong influence on the open multimodal chat ecosystem.",
        ],
        fill_rgb=PALE_GREEN,
        title_rgb=GREEN,
    )
    add_card(
        slide,
        Inches(8.78),
        Inches(1.05),
        Inches(3.88),
        Inches(2.35),
        "Main Limitations",
        [
            "Still depends on a non-trivial multi-GPU budget for faithful reproduction.",
            "Behavior is strongly tied to the quality of the generated instruction data pipeline.",
            "Outputs are mostly text, so it is not a dense vision generalist.",
        ],
        fill_rgb=PALE_RED,
        title_rgb=RED,
    )
    add_card(
        slide,
        Inches(0.85),
        Inches(3.85),
        Inches(11.75),
        Inches(1.72),
        "Best Way To Position LLaVA In This Deck",
        [
            "LLaVA is the bridge between pure multimodal alignment and practical multimodal assistants.",
            "It is easier to explain to users than BLIP-2, but less attractive for this particular reproduction because the main scientific interest here was staged bridge learning under hardware limits.",
        ],
        fill_rgb=CARD,
    )
    slide_no += 1

    slide = new_slide(prs, "VisionLLM v2 In A Nutshell", slide_no, "Sources: VisionLLM v2 paper, term_paper.tex")
    add_card(slide, Inches(0.72), Inches(1.3), Inches(1.7), Inches(1.2), "Image + Prompt", ["The user asks for a task in natural language."], font_size=13)
    add_card(slide, Inches(2.78), Inches(1.1), Inches(3.05), Inches(1.6), "Central MLLM", ["A shared multimodal backbone reasons over image and text."], fill_rgb=PALE_BLUE, font_size=15)
    add_card(slide, Inches(6.18), Inches(1.1), Inches(2.1), Inches(1.6), "Super Link", ["Routes task information and gradients between the MLLM and downstream decoders."], fill_rgb=PALE_GOLD, title_rgb=ACCENT, font_size=13)
    add_arrow(slide, Inches(2.4), Inches(1.68), Inches(0.26), Inches(0.22), HEADER)
    add_arrow(slide, Inches(5.84), Inches(1.68), Inches(0.26), Inches(0.22), ACCENT)
    add_card(slide, Inches(8.7), Inches(0.98), Inches(1.88), Inches(1.0), "Detection", ["Boxes / localization"], fill_rgb=PALE_GREEN, title_rgb=GREEN, font_size=12)
    add_card(slide, Inches(10.7), Inches(0.98), Inches(1.88), Inches(1.0), "Pose", ["Keypoints / structure"], fill_rgb=PALE_GREEN, title_rgb=GREEN, font_size=12)
    add_card(slide, Inches(8.7), Inches(2.12), Inches(1.88), Inches(1.0), "Generation", ["Image creation"], fill_rgb=PALE_GREEN, title_rgb=GREEN, font_size=12)
    add_card(slide, Inches(10.7), Inches(2.12), Inches(1.88), Inches(1.0), "Editing", ["Image manipulation"], fill_rgb=PALE_GREEN, title_rgb=GREEN, font_size=12)
    add_card(
        slide,
        Inches(0.82),
        Inches(3.65),
        Inches(5.7),
        Inches(2.0),
        "What It Is Really Trying To Solve",
        [
            "Most MLLMs only answer in text. VisionLLM v2 tries to break that limit by supporting structured and visual outputs too.",
            "It wants one model family to cover many task types that are usually handled by separate systems.",
        ],
        fill_rgb=CARD,
    )
    add_card(
        slide,
        Inches(6.8),
        Inches(3.65),
        Inches(5.7),
        Inches(2.0),
        "How It Differs From Tool-Using MLLMs",
        [
            "Tool-using MLLMs often pass text to external tools.",
            "VisionLLM v2 tries to connect the decoders more deeply so task information and gradients can move end to end.",
        ],
        fill_rgb=PALE_BLUE,
    )
    slide_no += 1

    slide = new_slide(prs, "VisionLLM v2 Method, Strengths, Limits", slide_no, "Sources: VisionLLM v2 paper, term_paper.tex")
    add_card(
        slide,
        Inches(0.7),
        Inches(1.0),
        Inches(3.85),
        Inches(2.6),
        "Method",
        [
            "Collects and harmonizes data from hundreds of vision and vision-language tasks.",
            "Joint-trains a shared MLLM with multiple task-specific decoders.",
            "Uses prompts to invoke different task behaviors through one shared system.",
        ],
        fill_rgb=PALE_GOLD,
        title_rgb=ACCENT,
    )
    add_card(
        slide,
        Inches(4.74),
        Inches(1.0),
        Inches(3.85),
        Inches(2.6),
        "Where It Is Good",
        [
            "Unifying perception, understanding, and generation.",
            "Supporting dense outputs such as localization and pose, not only text.",
            "Making the strongest 'generalist' claim among the three papers.",
        ],
        fill_rgb=PALE_GREEN,
        title_rgb=GREEN,
    )
    add_card(
        slide,
        Inches(8.78),
        Inches(1.0),
        Inches(3.85),
        Inches(2.6),
        "Main Limitations",
        [
            "Highest engineering complexity by far: many decoders, many datasets, many interfaces.",
            "Hard to reproduce, hard to ablate, and hard to shrink cleanly for a class project.",
            "Its ambition is impressive, but that same ambition is the reproducibility barrier.",
        ],
        fill_rgb=PALE_RED,
        title_rgb=RED,
    )
    add_card(
        slide,
        Inches(0.85),
        Inches(4.03),
        Inches(11.75),
        Inches(1.55),
        "Best Way To Position VisionLLM v2 In This Deck",
        [
            "It is the endpoint of the progression: the paper that broadens multimodal systems from text-answering models into multi-output vision generalists.",
            "That is why it is also the least realistic reproduction target under local student hardware.",
        ],
        fill_rgb=CARD,
    )
    slide_no += 1

    slide = new_slide(prs, "Main Takeaway From The Three-Paper Critique", slide_no, "Source: paper/term_paper.tex")
    add_progression_box(
        slide,
        Inches(0.72),
        Inches(1.3),
        Inches(3.78),
        Inches(1.8),
        "BLIP-2",
        "Ask whether a small trainable bridge can unlock frozen vision and language backbones.",
        HEADER,
    )
    add_progression_box(
        slide,
        Inches(4.58),
        Inches(1.3),
        Inches(3.78),
        Inches(1.8),
        "LLaVA",
        "Ask whether visual instruction tuning can make the model behave like a helpful assistant.",
        ACCENT,
    )
    add_progression_box(
        slide,
        Inches(8.44),
        Inches(1.3),
        Inches(3.85),
        Inches(1.8),
        "VisionLLM v2",
        "Ask whether one multimodal system can cover many output spaces and task families.",
        GREEN,
    )
    add_card(
        slide,
        Inches(0.9),
        Inches(3.75),
        Inches(11.55),
        Inches(2.0),
        "Overall Pattern",
        [
            "Capability rises from modular bridge learning to assistant behavior to full multimodal generality.",
            "At the same time, the data burden, engineering burden, and compute burden rise sharply.",
            "That tradeoff is the key reason BLIP-2 became the implementation target and the lens for the reproduction study.",
        ],
        fill_rgb=PALE_BLUE,
    )
    slide_no += 1

    slide = new_slide(prs, "Our Reproduction Target And Constraints", slide_no, "Sources: docs/blip2/README.md, metrics/blip2/env_snapshot.json")
    add_card(
        slide,
        Inches(0.78),
        Inches(1.1),
        Inches(5.45),
        Inches(2.5),
        "Target Of Reproduction",
        [
            "Run all three BLIP-2 stages end to end in a local environment.",
            "Produce checkpoints, caption outputs, and reproducible evaluation artifacts.",
            "Study the gap between pipeline-level success and paper-level performance.",
        ],
        fill_rgb=CARD,
    )
    env_rows = [
        ["Local setup", "Value"],
        ["OS", "Windows 11 Education"],
        ["Python", "3.10.11"],
        ["PyTorch", "2.4.1 + cu121"],
        ["GPU", "1x RTX 3070 8GB"],
        ["Codebase", "Official LAVIS"],
        ["Backbones", "CLIP ViT-L + OPT-350M"],
        ["Image size", "224"],
    ]
    add_table(
        slide,
        Inches(6.5),
        Inches(1.08),
        Inches(5.72),
        Inches(3.65),
        env_rows,
        col_widths=[Inches(2.15), Inches(3.57)],
        row_fills=[CARD, PALE_BLUE],
        font_size=12,
    )
    add_card(
        slide,
        Inches(0.78),
        Inches(4.15),
        Inches(11.45),
        Inches(1.7),
        "Important Scope Boundary",
        [
            "This was never a full paper-scale replication: the published BLIP-2 setup uses much larger backbones, much larger pretraining data, and much larger compute.",
            "The correct question was whether the method still behaves coherently after aggressive downscaling.",
        ],
        fill_rgb=PALE_GOLD,
        title_rgb=ACCENT,
    )
    slide_no += 1

    slide = new_slide(prs, "How We Implemented The Local BLIP-2 Reproduction", slide_no, "Sources: term_paper.tex, experiment_ledger.md, report_metrics_snapshot.json")
    add_card(slide, Inches(0.7), Inches(1.02), Inches(3.9), Inches(1.55), "Stage 1", ["Image-text representation learning", "Local schedules: 1 epoch, 3 epochs, or 10 epochs"], fill_rgb=PALE_BLUE)
    add_card(slide, Inches(4.72), Inches(1.02), Inches(3.9), Inches(1.55), "Stage 2", ["Generative alignment to OPT-350M", "Local schedules: 1 epoch, 3 epochs, or 10 epochs"], fill_rgb=PALE_GOLD, title_rgb=ACCENT)
    add_card(slide, Inches(8.74), Inches(1.02), Inches(3.9), Inches(1.55), "Caption Fine-Tuning", ["COCO caption adaptation", "Local schedules: 1, 5, or 15 epochs"], fill_rgb=PALE_GREEN, title_rgb=GREEN)
    rows = [
        ["Dimension", "Paper setup", "Local setup"],
        ["Vision backbone", "ViT-g", "CLIP ViT-L"],
        ["Language model", "OPT-2.7B", "OPT-350M"],
        ["Data scale", "Large multimodal mixture", "Reduced COCO Karpathy subsets"],
        ["Training budget", "250k + 80k steps", "1/1/1 to 10/10/15 schedules"],
        ["Evaluation", "Paper COCO caption table", "Offline subset-matched validation metrics"],
    ]
    add_table(
        slide,
        Inches(0.72),
        Inches(3.0),
        Inches(12.0),
        Inches(2.55),
        rows,
        col_widths=[Inches(2.25), Inches(4.1), Inches(5.65)],
        row_fills=[CARD, PALE_BLUE],
        font_size=12,
    )
    slide_no += 1

    slide = new_slide(prs, "Engineering Work That Was Actually Required", slide_no, "Sources: docs/blip2/failures_and_fixes.md, docs/blip2/README.md")
    cards = [
        ("Data path fixes", ["Absolute local annotation and image paths replaced broken default cache resolution."], CARD, HEADER),
        ("Single-GPU guards", ["Patched distributed assumptions so stage 1 and runners could work in a non-distributed local run."], PALE_BLUE, HEADER),
        ("Caption dataset fix", ["Stopped caption data image_id lists from being treated like retrieval tensors."], CARD, HEADER),
        ("OPT-350M bridge fix", ["Projected into OPT's actual input embedding dimension so stage 2 could train correctly."], PALE_GOLD, ACCENT),
        ("Evaluation rebuild", ["Saved predictions during training, then scored them offline against subset-matched COCO ground truth."], PALE_GREEN, GREEN),
        ("Windows workarounds", ["Handled METEOR instability, Java path quirks, optional spacy loading, and dependency drift."], PALE_RED, RED),
    ]
    x_positions = [Inches(0.72), Inches(4.42), Inches(8.12)]
    y_positions = [Inches(1.15), Inches(3.25)]
    idx = 0
    for row_y in y_positions:
        for col_x in x_positions:
            title, body, fill_rgb, title_rgb = cards[idx]
            add_card(slide, col_x, row_y, Inches(3.08), Inches(1.65), title, body, fill_rgb=fill_rgb, title_rgb=title_rgb, font_size=13, title_size=16)
            idx += 1
    add_card(
        slide,
        Inches(0.92),
        Inches(5.38),
        Inches(11.35),
        Inches(0.85),
        "Why This Matters",
        [
            "The project was not just a model run. It was a systems engineering effort to make a research codebase operate reproducibly on a single Windows workstation.",
        ],
        fill_rgb=CARD,
        font_size=14,
        title_size=16,
    )
    slide_no += 1

    local_categories = ["10k baseline", "50k best", "50k polish", "100k best", "100k final"]
    local_bleu4 = [
        REPORT["local_baseline_10k"]["metrics_x100"]["Bleu_4"],
        REPORT["local_office_best_50k"]["metrics_x100"]["Bleu_4"],
        REPORT["local_office_polish_final_50k"]["metrics_x100"]["Bleu_4"],
        round(RUN_100K_E2["metrics"]["Bleu_4"] * 100, 2),
        round(RUN_100K_E14["metrics"]["Bleu_4"] * 100, 2),
    ]
    local_cider = [
        REPORT["local_baseline_10k"]["metrics_x100"]["CIDEr"],
        REPORT["local_office_best_50k"]["metrics_x100"]["CIDEr"],
        REPORT["local_office_polish_final_50k"]["metrics_x100"]["CIDEr"],
        round(RUN_100K_E2["metrics"]["CIDEr"] * 100, 2),
        round(RUN_100K_E14["metrics"]["CIDEr"] * 100, 2),
    ]
    local_unique = [
        REPORT["local_baseline_10k"]["diversity"]["unique_caption_count"],
        REPORT["local_office_best_50k"]["diversity"]["unique_caption_count"],
        607,
        RUN_100K_E2["diversity"]["unique_caption_count"],
        RUN_100K_E14["diversity"]["unique_caption_count"],
    ]

    slide = new_slide(prs, "Experiment Family And Local Performance", slide_no, "Sources: report_metrics_snapshot.json, 100k evaluation summaries")
    add_chart(
        slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.65),
        Inches(1.1),
        Inches(7.0),
        Inches(4.35),
        local_categories,
        [
            ("BLEU-4", local_bleu4, ACCENT),
            ("CIDEr", local_cider, HEADER),
        ],
        y_max=42,
        legend=True,
    )
    add_chart(
        slide,
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(7.95),
        Inches(1.1),
        Inches(4.55),
        Inches(2.3),
        local_categories,
        [("Unique captions", local_unique, GREEN)],
        y_max=800,
        legend=False,
    )
    rows = [
        ["Run", "Schedule", "BLEU-4", "CIDEr", "Unique"],
        ["10k", "1/1/1", "1.45", "3.03", "61"],
        ["50k best", "3/3/5", "11.13", "37.57", "563"],
        ["50k polish", "caption only", "10.49", "37.14", "607"],
        ["100k best", "10/10/15", "8.57", "25.34", "677"],
        ["100k final", "10/10/15", "6.61", "18.36", "734"],
    ]
    add_table(
        slide,
        Inches(7.88),
        Inches(3.7),
        Inches(4.62),
        Inches(2.15),
        rows,
        col_widths=[Inches(1.32), Inches(1.05), Inches(0.83), Inches(0.82), Inches(0.6)],
        row_fills=[CARD, PALE_BLUE],
        font_size=10,
    )
    slide_no += 1

    slide = new_slide(prs, "Compared With The BLIP-2 Paper", slide_no, "Sources: BLIP-2 paper, term_paper.tex, report_metrics_snapshot.json")
    add_chart(
        slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.85),
        Inches(1.18),
        Inches(6.05),
        Inches(4.55),
        ["BLEU-4", "CIDEr"],
        [
            ("Paper: ViT-g OPT-2.7B", [43.7, 145.8], HEADER),
            ("Local best: 50k office epoch 3", [11.13, 37.57], ACCENT),
        ],
        y_max=160,
        legend=True,
    )
    add_metric_box(slide, Inches(7.35), Inches(1.3), Inches(2.2), Inches(1.0), "Paper BLEU-4", "43.7", PALE_BLUE)
    add_metric_box(slide, Inches(9.8), Inches(1.3), Inches(2.2), Inches(1.0), "Local BLEU-4", "11.13", PALE_GOLD)
    add_metric_box(slide, Inches(7.35), Inches(2.52), Inches(2.2), Inches(1.0), "Paper CIDEr", "145.8", PALE_BLUE)
    add_metric_box(slide, Inches(9.8), Inches(2.52), Inches(2.2), Inches(1.0), "Local CIDEr", "37.57", PALE_GOLD)
    add_card(
        slide,
        Inches(7.25),
        Inches(3.85),
        Inches(4.9),
        Inches(1.75),
        "Interpretation",
        [
            "The local result is far below the paper, but the comparison is not apples to apples.",
            "The local system uses smaller backbones, much less data, a reduced validation subset, and a far smaller compute budget.",
            "So this slide supports a scale-gap argument, not a leaderboard claim.",
        ],
        fill_rgb=CARD,
    )
    slide_no += 1

    slide = new_slide(prs, "Why The Gap Exists", slide_no, "Sources: term_paper.tex, failures_and_fixes.md")
    gap_cards = [
        ("Backbone mismatch", ["Paper: ViT-g + OPT-2.7B", "Local: CLIP ViT-L + OPT-350M"], PALE_BLUE, HEADER),
        ("Data mismatch", ["Paper uses massive multimodal pretraining.", "Local runs use reduced COCO subsets only."], PALE_GOLD, ACCENT),
        ("Optimization budget", ["Paper-scale steps and hardware were impossible locally.", "Student schedules were much shorter and cheaper."], CARD, HEADER),
        ("Evaluation mismatch", ["Local scores come from a consistent 1k validation subset, not the paper's full Karpathy test setup."], PALE_GREEN, GREEN),
        ("Environment friction", ["Windows-specific dependency and metric issues forced extra engineering and offline evaluation."], PALE_RED, RED),
    ]
    positions = [
        (Inches(0.82), Inches(1.18)),
        (Inches(4.34), Inches(1.18)),
        (Inches(7.86), Inches(1.18)),
        (Inches(2.6), Inches(3.35)),
        (Inches(7.15), Inches(3.35)),
    ]
    sizes = [
        (Inches(3.05), Inches(1.6)),
        (Inches(3.05), Inches(1.6)),
        (Inches(3.05), Inches(1.6)),
        (Inches(3.9), Inches(1.65)),
        (Inches(3.9), Inches(1.65)),
    ]
    for (title, body, fill_rgb, title_rgb), (x, y), (w, h) in zip(gap_cards, positions, sizes):
        add_card(slide, x, y, w, h, title, body, fill_rgb=fill_rgb, title_rgb=title_rgb, font_size=13, title_size=16)
    add_card(
        slide,
        Inches(0.95),
        Inches(5.42),
        Inches(11.3),
        Inches(0.82),
        "Bottom Line",
        [
            "The performance gap is the expected result of a faithful but aggressively downscaled reproduction. It does not mean the pipeline was broken.",
        ],
        fill_rgb=CARD,
        font_size=14,
        title_size=16,
    )
    slide_no += 1

    slide = new_slide(prs, "Qualitative Behavior Of The Reproduced Model", slide_no, "Sources: caption_eval_examples_office_epoch3.json, staged COCO images")
    add_card(
        slide,
        Inches(0.78),
        Inches(0.92),
        Inches(11.85),
        Inches(0.82),
        "Common Error Pattern",
        [
            "Predictions often sound fluent but substitute the wrong object, person, or scene. That is why diversity improved faster than true semantic accuracy.",
        ],
        fill_rgb=PALE_BLUE,
        font_size=14,
        title_size=16,
    )
    example_indices = [1, 2, 4]
    image_paths = [
        ROOT / "repo_study/LAVIS/cache/coco/images/val2014/COCO_val2014_000000511622.jpg",
        ROOT / "repo_study/LAVIS/cache/coco/images/val2014/COCO_val2014_000000341113.jpg",
        ROOT / "repo_study/LAVIS/cache/coco/images/val2014/COCO_val2014_000000353027.jpg",
    ]
    x_vals = [Inches(0.8), Inches(4.45), Inches(8.1)]
    for i, example_idx in enumerate(example_indices):
        ex = QUAL_EXAMPLES[example_idx]
        frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x_vals[i], Inches(1.98), Inches(3.05), Inches(4.45))
        frame.fill.solid()
        frame.fill.fore_color.rgb = CARD
        frame.line.color.rgb = LINE
        add_picture_contain(slide, image_paths[i], x_vals[i] + Inches(0.08), Inches(2.06), Inches(2.89), Inches(1.88))
        add_textbox(slide, x_vals[i] + Inches(0.09), Inches(4.03), Inches(2.86), Inches(0.78), [f"Reference: {ex['references'][0]}"], font_size=11, color=TEXT)
        add_textbox(slide, x_vals[i] + Inches(0.09), Inches(4.84), Inches(2.86), Inches(1.28), [f"Prediction: {ex['prediction']}"], font_size=11, color=RED)
    slide_no += 1

    slide = new_slide(prs, "What This Project Actually Proved", slide_no, "Sources: term_paper.tex, mkamrul_project_answers.md")
    add_card(
        slide,
        Inches(0.9),
        Inches(1.08),
        Inches(5.45),
        Inches(3.55),
        "Succeeded",
        [
            "The full BLIP-2 pipeline ran locally: stage 1, stage 2, and caption fine-tuning all completed.",
            "The reproduction produced reusable scripts, checkpoints, logs, metric files, and evaluation artifacts.",
            "Caption quality improved in an interpretable way from 10k to the 50k office run.",
        ],
        fill_rgb=PALE_GREEN,
        title_rgb=GREEN,
    )
    add_card(
        slide,
        Inches(6.95),
        Inches(1.08),
        Inches(5.45),
        Inches(3.55),
        "Did Not Match",
        [
            "The local model stayed far below the paper's captioning quality.",
            "The 100k long run increased diversity but did not beat the earlier 50k best on BLEU-4 or CIDEr.",
            "So pipeline reproducibility did not imply performance reproducibility.",
        ],
        fill_rgb=PALE_RED,
        title_rgb=RED,
    )
    add_card(
        slide,
        Inches(1.15),
        Inches(5.02),
        Inches(11.0),
        Inches(0.9),
        "Final Thesis",
        [
            "BLIP-2 was the correct reproduction target because its core method survives downscaling. The project's real result is that modular multimodal pipelines can be reproduced on student hardware even when paper-level performance cannot.",
        ],
        fill_rgb=CARD,
        font_size=14,
        title_size=16,
    )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))


if __name__ == "__main__":
    build_presentation()
